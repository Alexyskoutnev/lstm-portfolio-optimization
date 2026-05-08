"""Build a PowerPoint deck for the four-arm portfolio benchmark project.

Reads the existing presentation_*.png plots in plots/ and assembles them
into a 16:9 PPTX deck with a title page, RQ-by-RQ result pages, methodology,
limitations, and conclusion. Output: docs/presentation/portfolio_ladder_deck.pptx

The deck inherits design choices from our matplotlib presentation theme:
- Same per-arm color palette (Okabe-Ito) referenced in titles/accents
- Sans-serif, large titles, neutral background
- One image per slide, with concise bullet text in the speaker-notes pane

Run: ``python scripts/build_presentation_deck.py``
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
PLOTS_DIR = PROJECT_ROOT / "plots"
COVERS_DIR = PLOTS_DIR / "covers"
OUTPUT_DIR = PROJECT_ROOT / "docs" / "presentation"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_PATH = OUTPUT_DIR / "portfolio_ladder_deck.pptx"

# 16:9 dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Visual identity — same hex codes used in the matplotlib plots
NAVY = RGBColor(0x0E, 0x1A, 0x2B)
WHITE = RGBColor(0xFA, 0xFA, 0xFA)
SOFT_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT_GREEN = RGBColor(0x00, 0x9E, 0x73)   # GBT
ACCENT_BLUE = RGBColor(0x00, 0x72, 0xB2)    # LSTM
ACCENT_ORANGE = RGBColor(0xE6, 0x9F, 0x00)  # Ridge
NEUTRAL_GRAY = RGBColor(0x88, 0x88, 0x88)   # Sample mean


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _solid_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide,
    left: Emu, top: Emu, width: Emu, height: Emu,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = NAVY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    font.name = "Helvetica"
    return tf


def _add_bullets(
    slide,
    left: Emu, top: Emu, width: Emu, height: Emu,
    items: list[str],
    *,
    size: int = 20,
    color: RGBColor = NAVY,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + item
        font = run.font
        font.size = Pt(size)
        font.color.rgb = color
        font.name = "Helvetica"


def _add_accent_bar(slide, color: RGBColor) -> None:
    """Small left-edge color bar — visual identity carried across slides."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.18), SLIDE_H,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _add_footer(slide, page_number: int, total: int) -> None:
    _add_text(
        slide,
        Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
        "AIM 5005 Final Project · Risk-Adjusted Portfolio Optimization · "
        "Skoutnev, Mukahanana, Chiremba",
        size=10, color=SOFT_GRAY,
    )
    _add_text(
        slide,
        Inches(11.5), Inches(7.05), Inches(1.7), Inches(0.35),
        f"{page_number} / {total}",
        size=10, color=SOFT_GRAY, align=PP_ALIGN.RIGHT,
    )


def _add_title(
    slide, title: str, subtitle: str | None = None, accent: RGBColor = NAVY
) -> None:
    _add_text(
        slide,
        Inches(0.6), Inches(0.35), Inches(12.2), Inches(0.7),
        title, size=32, bold=True, color=accent,
    )
    if subtitle:
        _add_text(
            slide,
            Inches(0.6), Inches(1.05), Inches(12.2), Inches(0.4),
            subtitle, size=15, color=SOFT_GRAY,
        )


def _add_image_centered(
    slide, image_path: Path,
    *, top: float = 1.6, max_height: float = 5.2, max_width: float = 12.0,
) -> None:
    """Add an image centered in the lower 2/3 of the slide."""
    if not image_path.exists():
        return
    pic = slide.shapes.add_picture(
        str(image_path), Inches(0), Inches(top),
        height=Inches(max_height),
    )
    # Constrain by width if too wide
    if pic.width > Inches(max_width):
        pic.width = Inches(max_width)
        pic.height = Inches(max_width * pic.height.inches / pic.width.inches)
    # Center horizontally
    pic.left = Inches((13.333 - pic.width.inches) / 2)


def _add_full_bleed_background(slide, image_path: Path) -> None:
    """Place an image full-bleed behind any other content on the slide."""
    if not image_path.exists():
        return
    pic = slide.shapes.add_picture(
        str(image_path), Inches(0), Inches(0),
        width=SLIDE_W, height=SLIDE_H,
    )
    # Move to back
    spTree = pic._element.getparent()
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)


def _add_dim_overlay(
    slide, *, top: float, height: float,
    color: RGBColor = NAVY, opacity_pct: int = 70,
) -> None:
    """Translucent rectangle overlay for text legibility on photo backgrounds."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(top), SLIDE_W, Inches(height),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    # python-pptx doesn't expose alpha directly via the public API; use the
    # underlying XML to set transparency on the fill.
    from pptx.oxml.ns import qn
    sp = rect.fill._xPr
    solid_fill = sp.find(qn("a:solidFill"))
    if solid_fill is not None:
        srgb = solid_fill.find(qn("a:srgbClr"))
        if srgb is not None:
            from lxml import etree
            alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", str(opacity_pct * 1000))


# ---------------------------------------------------------------------------
# Slide content
# ---------------------------------------------------------------------------

def slide_title(prs: Presentation, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, NAVY)
    _add_full_bleed_background(s, COVERS_DIR / "title_cover.png")
    # Dark scrim across the top half so title is legible regardless of image
    _add_dim_overlay(s, top=0.0, height=4.5, color=NAVY, opacity_pct=55)
    # Big title
    _add_text(
        s, Inches(1.0), Inches(0.9), Inches(11.5), Inches(1.8),
        "Risk-Adjusted Portfolio Optimization",
        size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
    )
    _add_text(
        s, Inches(1.0), Inches(2.5), Inches(11.5), Inches(0.7),
        "ML vs. Markowitz — A Four-Arm Model-Complexity Ladder",
        size=24, color=ACCENT_GREEN, align=PP_ALIGN.LEFT,
    )
    # Bottom credits with their own scrim
    _add_dim_overlay(s, top=6.3, height=1.2, color=NAVY, opacity_pct=70)
    _add_text(
        s, Inches(1.0), Inches(6.45), Inches(11.5), Inches(0.4),
        "Alexy Skoutnev  ·  Shaun Mukahanana  ·  Tadiwanashe Chiremba",
        size=16, color=WHITE, align=PP_ALIGN.LEFT,
    )
    _add_text(
        s, Inches(1.0), Inches(6.85), Inches(11.5), Inches(0.4),
        "AIM 5005 — Yeshiva University — May 2026",
        size=13, color=SOFT_GRAY, align=PP_ALIGN.LEFT,
    )


def _section_divider(
    prs: Presentation,
    page: int, total: int,
    *,
    cover_filename: str,
    eyebrow: str,
    title: str,
    subtitle: str | None = None,
    accent: RGBColor = ACCENT_GREEN,
) -> None:
    """Full-bleed AI-image section divider with text overlay on a dark scrim."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, NAVY)
    _add_full_bleed_background(s, COVERS_DIR / cover_filename)
    # Lower-left scrim for legibility
    _add_dim_overlay(s, top=4.0, height=3.5, color=NAVY, opacity_pct=68)
    _add_text(
        s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.5),
        eyebrow.upper(), size=14, bold=True, color=accent,
        align=PP_ALIGN.LEFT,
    )
    _add_text(
        s, Inches(0.8), Inches(4.85), Inches(11.7), Inches(1.4),
        title, size=44, bold=True, color=WHITE,
        align=PP_ALIGN.LEFT,
    )
    if subtitle:
        _add_text(
            s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6),
            subtitle, size=18, color=SOFT_GRAY,
            align=PP_ALIGN.LEFT,
        )
    _add_footer(s, page, total)


def slide_section_methodology(prs: Presentation, page: int, total: int) -> None:
    _section_divider(
        prs, page, total,
        cover_filename="method_concept.png",
        eyebrow="Section",
        title="Methodology",
        subtitle="Same data, same backtest, same Σ — only µ differs",
        accent=ACCENT_GREEN,
    )


def slide_section_results(prs: Presentation, page: int, total: int) -> None:
    _section_divider(
        prs, page, total,
        cover_filename="rq1_hero.png",
        eyebrow="Section",
        title="Results",
        subtitle="Test period 2022 → 2025  ·  rolling-window walk-forward backtest",
        accent=ACCENT_BLUE,
    )


def slide_section_risk(prs: Presentation, page: int, total: int) -> None:
    _section_divider(
        prs, page, total,
        cover_filename="rq3_hero.png",
        eyebrow="Section",
        title="The Risk Picture",
        subtitle="Drawdown, VaR, and CVaR — does ML buy us tail-protection?",
        accent=ACCENT_BLUE,
    )


def slide_section_reflections(prs: Presentation, page: int, total: int) -> None:
    _section_divider(
        prs, page, total,
        cover_filename="tuning_concept.png",
        eyebrow="Section",
        title="What We Learned",
        subtitle="Tuning choices that flipped the result, and what we'd do next",
        accent=ACCENT_ORANGE,
    )


def slide_question(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_GREEN)
    _add_title(s, "The Research Question",
               "Three questions from the project proposal")
    _add_bullets(
        s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.0),
        [
            "RQ1: Can ML-based forecasts improve Sharpe / Sortino vs classical MVO?",
            "RQ2: Are ML-enhanced portfolios more robust during high-volatility regimes?",
            "RQ3: Does nonlinear temporal structure reduce downside risk?",
            "",
            "We evaluate ML vs MVO on 11 GICS sector ETFs, 2010–2025,",
            "with rolling-window backtests on a held-out 2022–2025 test slice.",
        ],
        size=22,
    )
    _add_footer(s, page, total)


def slide_approach(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_GREEN)
    _add_title(s, "Our Approach: A Complexity Ladder",
               "Same Markowitz pipeline; only the µ estimator changes")
    _add_bullets(
        s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.0),
        [
            "1. Sample mean    — no learning (baseline)",
            "2. Ridge          — linear ML on a tabular feature panel",
            "3. LightGBM       — gradient-boosted trees on the same panel",
            "4. LSTM           — recurrent net on multi-channel sequences",
            "",
            "All four forecast the same 21-day forward return per asset.",
            "Σ stays sample-based.  Optimizer stays SLSQP, long-only.  Backtest stays identical.",
        ],
        size=22,
    )
    _add_footer(s, page, total)


def slide_pipeline(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_GREEN)
    _add_title(s, "Architecture: Only µ Differs Between Arms")
    _add_image_centered(s, PLOTS_DIR / "mvo_vs_gbt_pipeline.png",
                        top=1.4, max_height=5.5)
    _add_footer(s, page, total)


def slide_mu_zoom(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_GREEN)
    _add_title(s, "The µ Swap, Zoomed",
               "The only step that changes between classical MVO and any ML arm")
    _add_image_centered(s, PLOTS_DIR / "mvo_vs_gbt_mu_zoom.png",
                        top=1.7, max_height=4.8)
    _add_footer(s, page, total)


def slide_methodology(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_GREEN)
    _add_title(s, "Methodology",
               "Walk-forward backtest with rigorous no-look-ahead discipline")
    _add_bullets(
        s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.0),
        [
            "Universe: 11 S&P 500 sector ETFs (XLK, XLF, XLE, …, XLC)",
            "Period: 2010-01-01 → 2025-12-31  (~4,000 trading days)",
            "Split: 70 / 15 / 15 temporal — train ends 2021-03",
            "Target: cumulative log return over the next 21 trading days",
            "Features: lagged returns, rolling vol, VIX features, regime, cross-asset, calendar",
            "Backtest: 252-day lookback, 21-day rebalance, max-Sharpe, long-only",
            "Models trained on train slice only.  Test slice: 2022-03 → 2025-12.",
        ],
        size=20,
    )
    _add_footer(s, page, total)


def slide_hero(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_GREEN)
    _add_title(s, "Headline: GBT Wins",
               "All four arms · same Σ, optimizer, constraints · only µ differs")
    _add_image_centered(s, PLOTS_DIR / "presentation_hero.png",
                        top=1.4, max_height=5.5)
    _add_footer(s, page, total)


def slide_equity(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_GREEN)
    _add_title(s, "Equity Curves — Test Period 2022–2025",
               "GBT 79% cumulative · Sample-mean 35% · LSTM 35% · Ridge 26%")
    _add_image_centered(s, PLOTS_DIR / "presentation_equity_curves.png",
                        top=1.5, max_height=5.4)
    _add_footer(s, page, total)


def slide_rq1(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_BLUE)
    _add_title(s, "RQ1 — Better Sharpe / Sortino than classical MVO?",
               "ANSWER: YES — GBT delivers Sharpe 1.46 vs 0.96 for sample-mean MVO (+53%)")
    _add_image_centered(s, PLOTS_DIR / "presentation_metrics_bars.png",
                        top=1.6, max_height=4.8)
    _add_text(
        s, Inches(0.6), Inches(6.55), Inches(12.0), Inches(0.5),
        "Ridge actively hurts (Sharpe 0.61).  ML lift comes from non-linearity, not from \"having a model.\"",
        size=14, color=SOFT_GRAY, align=PP_ALIGN.CENTER,
    )
    _add_footer(s, page, total)


def slide_rq2(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_BLUE)
    _add_title(s, "RQ2 — Robust in High-Vol Regimes?",
               "ANSWER: YES — GBT high-vol Sharpe is −1.1 vs −3.5 for sample-mean (3× less degradation)")
    _add_image_centered(s, PLOTS_DIR / "presentation_regime_sharpe.png",
                        top=1.6, max_height=4.8)
    _add_text(
        s, Inches(0.6), Inches(6.55), Inches(12.0), Inches(0.5),
        "GBT also wins in medium-vol (1.6 vs 0.2 baseline).  Non-linear conditioning on regime indicators pays off most in stress.",
        size=14, color=SOFT_GRAY, align=PP_ALIGN.CENTER,
    )
    _add_footer(s, page, total)


def slide_rq3(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_BLUE)
    _add_title(s, "RQ3 — Does Nonlinearity Reduce Downside Risk?",
               "ANSWER: YES — LSTM is safest on every metric (max DD, VaR, CVaR)")
    _add_image_centered(s, PLOTS_DIR / "presentation_var_cvar.png",
                        top=1.5, max_height=4.5)
    _add_text(
        s, Inches(0.6), Inches(6.30), Inches(12.0), Inches(0.4),
        "LSTM Max DD −12.05%  ·  VaR 95%  −1.13%  ·  CVaR 95%  −1.60%",
        size=18, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER,
    )
    _add_text(
        s, Inches(0.6), Inches(6.75), Inches(12.0), Inches(0.4),
        "LSTM trades return for safety.  GBT also beats the baseline on VaR.  Ridge is worst on every downside metric.",
        size=13, color=SOFT_GRAY, align=PP_ALIGN.CENTER,
    )
    _add_footer(s, page, total)


def slide_drawdown(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_BLUE)
    _add_title(s, "Drawdown — Underwater View",
               "Lower swings means tighter risk control")
    _add_image_centered(s, PLOTS_DIR / "presentation_drawdown.png",
                        top=1.7, max_height=5.0)
    _add_footer(s, page, total)


def slide_complexity(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_GREEN)
    _add_title(s, "The Complexity Ladder",
               "Sharpe and Sortino across the four-arm sequence")
    _add_image_centered(s, PLOTS_DIR / "presentation_complexity_ladder.png",
                        top=1.5, max_height=5.2)
    _add_text(
        s, Inches(0.6), Inches(6.85), Inches(12.0), Inches(0.4),
        "Non-monotonic: trees beat linear and recurrent.  Linearity is too rigid; full recurrence overfits raw return sequences.",
        size=13, color=SOFT_GRAY, align=PP_ALIGN.CENTER,
    )
    _add_footer(s, page, total)


def slide_features(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_GREEN)
    _add_title(s, "What is the GBT Looking At?",
               "Top-15 features by gain — model learns regime, momentum, dispersion")
    _add_image_centered(s, PLOTS_DIR / "presentation_gbt_importance.png",
                        top=1.5, max_height=5.4)
    _add_footer(s, page, total)


def slide_turnover(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_BLUE)
    _add_title(s, "Transaction-Cost Picture",
               "Average daily L1 weight change — proxy for trading-cost burden")
    _add_image_centered(s, PLOTS_DIR / "presentation_turnover.png",
                        top=1.6, max_height=4.8)
    _add_text(
        s, Inches(0.6), Inches(6.55), Inches(12.0), Inches(0.5),
        "LSTM is cheapest to trade.  GBT trades 3× as aggressively as sample-mean — its return advantage comes with real-world friction.",
        size=14, color=SOFT_GRAY, align=PP_ALIGN.CENTER,
    )
    _add_footer(s, page, total)


def slide_tuning_lessons(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_ORANGE)
    _add_title(s, "What Mattered — Tuning Lessons",
               "Three fixes that flipped \"ML doesn't help\" into a clean win")
    _add_bullets(
        s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.0),
        [
            "Drop validation-driven early stopping.  Our val window straddles the 2020 COVID crash, "
            "which made val loss explode and stopped GBT after 2 trees.  Fixed iteration + L2 reg fixed it.",
            "Standardize features for Ridge.  Without z-scoring, the L2 penalty falls almost entirely "
            "on small-scale return features — they're effectively unused.",
            "Multi-channel sequences + Huber loss for LSTM.  Single-channel return sequences are too "
            "sparse a signal; MSE rewards \"predict zero\" too much on fat-tailed targets.",
            "",
            "General lesson: \"ML didn't help\" on a financial task is almost always a tuning artifact "
            "before it's an empirical finding.  Validate carefully.",
        ],
        size=18,
    )
    _add_footer(s, page, total)


def slide_limitations(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, WHITE)
    _add_accent_bar(s, ACCENT_ORANGE)
    _add_title(s, "Limitations & Future Work")
    _add_bullets(
        s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.5),
        [
            "Single test period (2022–2025).  Walk-forward retraining across multiple test windows "
            "would build confidence the GBT lift isn't regime-specific.",
            "No transaction-cost model.  Turnover is a proxy; explicit cost simulation could change the "
            "GBT-vs-LSTM ranking.",
            "Σ stays sample-based.  Adding shrinkage (Ledoit-Wolf) or a learned Σ might further reduce "
            "Markowitz error amplification.",
            "Limited hyperparameter search.  Optuna sweeps for GBT and a richer LSTM architecture "
            "search are obvious next steps.",
            "No ensemble.  Combining GBT (return) and LSTM (defense) with a regime-aware blend would "
            "likely beat either alone.",
        ],
        size=18,
    )
    _add_footer(s, page, total)


def slide_conclusion(prs: Presentation, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _solid_background(s, NAVY)
    _add_full_bleed_background(s, COVERS_DIR / "conclusion_cover.png")
    # Top scrim for the takeaway
    _add_dim_overlay(s, top=0.0, height=3.0, color=NAVY, opacity_pct=70)
    _add_text(
        s, Inches(1.0), Inches(0.6), Inches(11.5), Inches(0.7),
        "Conclusion",
        size=44, bold=True, color=WHITE,
    )
    _add_text(
        s, Inches(1.0), Inches(1.6), Inches(11.5), Inches(0.7),
        "ML-forecasted µ improves Markowitz — when tuned correctly.",
        size=22, color=ACCENT_GREEN,
    )
    # Bottom scrim for the bullets
    _add_dim_overlay(s, top=3.6, height=3.9, color=NAVY, opacity_pct=72)
    _add_bullets(
        s, Inches(1.0), Inches(3.9), Inches(11.5), Inches(3.0),
        [
            "GBT-MVO: Sharpe 1.46 (+53% vs baseline), Sortino 2.14 (+67%)",
            "LSTM-MVO: lowest drawdown (-12.05%) and tightest tails (CVaR -1.60%)",
            "Ridge underperforms — the lift is non-linearity, not learning per se",
            "Largest relative ML win is in stressed regimes — exactly what RQ2 hypothesized",
        ],
        size=18, color=WHITE,
    )
    _add_text(
        s, Inches(1.0), Inches(6.85), Inches(11.5), Inches(0.5),
        "github.com/Alexyskoutnev/lstm-portfolio-optimization",
        size=13, color=SOFT_GRAY,
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_question,
        slide_approach,
        slide_pipeline,
        slide_mu_zoom,
        slide_section_methodology,    # AI hero divider
        slide_methodology,
        slide_section_results,        # AI hero divider
        slide_hero,
        slide_equity,
        slide_rq1,
        slide_rq2,
        slide_section_risk,           # AI hero divider
        slide_rq3,
        slide_drawdown,
        slide_complexity,
        slide_features,
        slide_turnover,
        slide_section_reflections,    # AI hero divider
        slide_tuning_lessons,
        slide_limitations,
        slide_conclusion,
    ]
    total = len(builders)
    for i, builder in enumerate(builders, start=1):
        if i == 1:
            builder(prs, total)
        else:
            builder(prs, i, total)

    prs.save(str(OUTPUT_PATH))
    print(f"wrote {OUTPUT_PATH}  ({total} slides)")


if __name__ == "__main__":
    main()
